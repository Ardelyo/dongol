#!/usr/bin/env python3
"""
Generator Presentasi PowerPoint DONGOL
Untuk presentasi makalah
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Buat presentasi
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Format title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0x20, 0x20, 0x60)
    
    return slide

def add_content_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    
    tf = body_shape.text_frame
    tf.clear()
    
    for i, text in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = text
        p.font.size = Pt(20)
        p.level = 0
    
    return slide

def add_image_slide(prs, title, image_path):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add image
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(11))
    
    return slide

# Slide 1: Title
add_title_slide(
    prs,
    "DONGOL",
    "Distributed Orchestration for Navigating Goals and Operational Logic\n\n"
    "Sistem Manajemen Tugas Paralel Modern\n\n"
    "Disusun oleh:\n"
    "Ardellio Satria Anindito\n"
    "SMA Kartika XIX-1 Bandung"
)

# Slide 2: Latar Belakang
add_content_slide(prs, "Latar Belakang", [
    "• Kebutuhan pemrosesan data efisien semakin meningkat",
    "• Sistem sekuensial mengalami kendala dengan data besar",
    "• Pendekatan \"Think Parallel, Execute Faster\"",
    "• DONGOL: Solusi manajemen tugas paralel dengan chunking cerdas",
    "• Agent-native design untuk AI workflows"
])

# Slide 3: Arsitektur
add_image_slide(prs, "Arsitektur Sistem DONGOL", "makalah_assets/architecture.png")

# Slide 4: Komponen Utama
add_content_slide(prs, "Komponen Utama", [
    "1. Presentation Layer - CLI, REST API, WebSocket, SDK",
    "2. Unified Core Engine - Async event loop, scheduler",
    "3. Parallel Thinker Matrix - 8 workers, load balancer",
    "4. Intelligent Chunking - Token, Structure, Semantic",
    "5. Context Memory - L1-L4 caching, Vector DB",
    "6. Persistence Layer - SQLite, Sled, File System"
])

# Slide 5: Teknologi
add_content_slide(prs, "Teknologi yang Digunakan", [
    "• Python 3.9+ dengan asyncio untuk concurrency",
    "• ThreadPool & ProcessPool untuk parallel execution",
    "• FastAPI untuk REST API",
    "• Rich & Click untuk CLI interaktif",
    "• Topological sort untuk dependency resolution",
    "• Actor-model dengan message passing"
])

# Slide 6: Hasil Benchmark
add_image_slide(prs, "Hasil Benchmarking", "makalah_assets/benchmark_graphs.png")

# Slide 7: Analisis Drive
add_image_slide(prs, "Analisis Drive D (Data Nyata)", "makalah_assets/file_distribution.png")

# Slide 8: Hasil Pengujian
add_content_slide(prs, "Hasil Pengujian Performa", [
    "• Task Creation: 48.713 tugas/detik (target: 10.000)",
    "• Parallel Speedup: 4.5x lebih cepat dari sekuensial",
    "• Throughput: 196 file/detik vs 43.7 file/detik",
    "• Analisis Drive: 2.979 file dalam 48 detik",
    "• File terorganisir: 1.455 file (5.36 GB)",
    "• Semua target performa terlampaui!"
])

# Slide 9: Perbandingan
add_content_slide(prs, "Perbandingan dengan Sistem Lain", [
    "Celery/RQ:",
    "  • Tidak ada parallel thinking",
    "  • Tidak ada intelligent chunking",
    "",
    "Airflow/Prefect:",
    "  • Latency dalam menit",
    "  • Manual chunking",
    "",
    "DONGOL:",
    "  • ✓ Native parallel thinking",
    "  • ✓ Automatic intelligent chunking",
    "  • ✓ Latency dalam milidetik"
])

# Slide 10: Kesimpulan
add_content_slide(prs, "Kesimpulan", [
    "1. DONGOL berhasil mengimplementasikan sistem paralel",
    "   dengan speedup 4.5x",
    "",
    "2. Intelligent chunking dengan dependency tracking",
    "   berfungsi optimal",
    "",
    "3. Task creation rate 48.713/s melampaui target",
    "",
    "4. Validasi pada data nyata (23.12 GB) berhasil",
    "",
    "5. Sistem siap untuk pengembangan lebih lanjut"
])

# Slide 11: Saran
add_content_slide(prs, "Saran Pengembangan", [
    "• Pengembangan GUI untuk aksesibilitas lebih baik",
    "• Implementasi fitur undo/redo",
    "• Mode terdistribusi multi-node",
    "• Integrasi cloud storage",
    "• Optimasi dengan Rust untuk hot paths"
])

# Slide 12: Terima Kasih
slide_layout = prs.slide_layouts[6]  # Blank
slide = prs.slides.add_slide(slide_layout)

# Add centered text
box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(3))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "Terima Kasih"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = RGBColor(0x20, 0x20, 0x60)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "\nArdellio Satria Anindito"
p2.font.size = Pt(28)
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = "SMA Kartika XIX-1 Bandung"
p3.font.size = Pt(24)
p3.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
p3.alignment = PP_ALIGN.CENTER

# Save
prs.save('Presentasi_DONGOL_Ardellio.pptx')
print("✓ Presentasi PowerPoint berhasil dibuat!")
print("  File: Presentasi_DONGOL_Ardellio.pptx")
print("  Total slide: 12")
